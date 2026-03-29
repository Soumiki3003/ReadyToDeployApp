import io
import logging
from pathlib import Path

import fitz
import pytesseract
from bs4 import BeautifulSoup
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation

from app import models, schemas


class FileService:
    def __init__(self):
        self.__logger = logging.getLogger(__name__)

    def extract_textual_content(self, filepath: Path) -> list[schemas.TextualContent]:
        self.__logger.info(f"Starting textual content extraction from {filepath}")
        text_data = []

        if filepath.suffix == ".pdf":
            self.__logger.debug(f"Processing PDF file: {filepath.name}")
            try:
                doc = fitz.open(filepath)
                self.__logger.debug(f"PDF opened successfully, total pages: {len(doc)}")
                for i, page in enumerate(doc.pages(), 1):
                    text = page.get_text("text").strip()
                    self.__logger.debug(
                        f"Extracted text from page {i}, length: {len(text)} chars"
                    )
                    text_data.append(
                        schemas.PaginatedTextualContent(
                            page=i, text=text, from_image=False
                        )
                    )
                self.__logger.info(
                    f"Successfully extracted text from {len(text_data)} pages"
                )
            except Exception as e:
                self.__logger.error(f"Error processing PDF {filepath.name}: {e}")
                raise

        elif filepath.suffix == ".pptx":
            self.__logger.debug(f"Processing PPTX file: {filepath.name}")
            try:
                prs = Presentation(filepath.as_posix())
                self.__logger.debug(
                    f"PPTX opened successfully, total slides: {len(prs.slides)}"
                )
                for i, slide in enumerate(prs.slides, 1):
                    slide_text = " ".join(
                        getattr(shape, "text")
                        for shape in slide.shapes
                        if hasattr(shape, "text")
                    )
                    self.__logger.debug(
                        f"Extracted text from slide {i}, length: {len(slide_text)} chars"
                    )
                    text_data.append(
                        schemas.SlideTextualContent(
                            slide=i, text=slide_text, from_image=False
                        )
                    )
                self.__logger.info(
                    f"Successfully extracted text from {len(text_data)} slides"
                )
            except Exception as e:
                self.__logger.error(f"Error processing PPTX {filepath.name}: {e}")
                raise

        elif filepath.suffix == ".html":
            self.__logger.debug(f"Processing HTML file: {filepath.name}")
            try:
                with open(filepath, "r", encoding="utf-8") as f:
                    soup = BeautifulSoup(f, "html.parser")
                body_text = soup.get_text(separator=" ", strip=True)
                self.__logger.debug(
                    f"Extracted HTML text, length: {len(body_text)} chars"
                )
                text_data.append(
                    schemas.HTMLTextualContent(
                        section="html", text=body_text, from_image=False
                    )
                )
                self.__logger.info("Successfully extracted text from HTML file")
            except Exception as e:
                self.__logger.error(f"Error processing HTML {filepath.name}: {e}")
                raise

        else:
            self.__logger.warning(
                f"Unsupported file type for textual content extraction: {filepath.suffix}"
            )

        self.__logger.info(
            f"Textual extraction complete. Total entries: {len(text_data)}"
        )
        return text_data

    def extract_visual_content(self, filepath: Path) -> list[schemas.TextualContent]:
        self.__logger.info(f"Starting visual content extraction from {filepath}")
        visual_text = []

        if filepath.suffix == ".pdf":
            self.__logger.debug(f"Processing PDF for OCR: {filepath.name}")
            try:
                self.__logger.debug("Converting PDF pages to images...")
                pages = convert_from_path(filepath, dpi=200)
                self.__logger.debug(
                    f"Successfully converted {len(pages)} pages to images"
                )
                for i, page_img in enumerate(pages, 1):
                    self.__logger.debug(f"Running OCR on page {i}...")
                    ocr_text = pytesseract.image_to_string(page_img)
                    self.__logger.debug(
                        f"OCR extracted text from page {i}, length: {len(ocr_text.strip())} chars"
                    )
                    visual_text.append(
                        schemas.PaginatedTextualContent(
                            page=i, text=ocr_text.strip(), from_image=True
                        )
                    )
                self.__logger.info(
                    f"Successfully extracted OCR text from {len(visual_text)} pages"
                )
            except Exception as e:
                error_msg = str(e)
                if "poppler" in error_msg.lower():
                    self.__logger.error("Poppler not found in PATH")
                    raise RuntimeError(
                        "Unable to get page count. Is poppler installed and in PATH?"
                    ) from e
                self.__logger.error(f"Error during PDF OCR processing: {e}")
                raise

        elif filepath.suffix == ".pptx":
            self.__logger.debug(f"Processing PPTX for OCR: {filepath.name}")
            try:
                prs = Presentation(filepath.as_posix())
                self.__logger.debug(
                    f"PPTX opened, processing {len(prs.slides)} slides for images"
                )
                for i, slide in enumerate(prs.slides, 1):
                    self.__logger.debug(f"Scanning slide {i} for images...")
                    image_count = 0
                    for shape in slide.shapes:
                        if shape.shape_type == 13:  # picture
                            image_count += 1
                            image = getattr(shape, "image", None)
                            if image is None:
                                self.__logger.debug(
                                    f"Slide {i}: Image {image_count} has no data, skipping"
                                )
                                continue
                            self.__logger.debug(
                                f"Slide {i}: Running OCR on image {image_count}"
                            )
                            image_bytes = io.BytesIO(image.blob)
                            img = Image.open(image_bytes)
                            ocr_text = pytesseract.image_to_string(img)
                            if not isinstance(ocr_text, str):
                                self.__logger.warning(
                                    f"OCR returned non-string result for slide {i} in {filepath.name}"
                                )
                                continue
                            ocr_text = ocr_text.strip()
                            self.__logger.debug(
                                f"Slide {i}: OCR extracted text, length: {len(ocr_text)} chars"
                            )
                            visual_text.append(
                                schemas.SlideTextualContent(
                                    slide=i,
                                    text=ocr_text.strip(),
                                    from_image=True,
                                )
                            )
                    self.__logger.debug(f"Slide {i}: Found {image_count} images")
                self.__logger.info(
                    f"Successfully extracted OCR text from {len(visual_text)} slide images"
                )
            except Exception as e:
                self.__logger.error(f"Error during PPTX OCR processing: {e}")
                raise
        else:
            self.__logger.warning(
                f"Unsupported file type for visual content extraction: {filepath.suffix}"
            )

        self.__logger.info(
            f"Visual extraction complete. Total entries: {len(visual_text)}"
        )
        return visual_text

    def chunk_textual_content(
        self,
        textual_contents: list[schemas.TextualContent],
        *,
        source_file: str,
        course_id: str,
    ) -> list[models.ContentChunk]:
        self.__logger.info(f"Chunking {len(textual_contents)} textual content entries")
        chunk_index = 0
        chunks: list[models.ContentChunk] = []

        for entry in textual_contents:
            if isinstance(entry, schemas.PaginatedTextualContent):
                page = entry.page
            elif isinstance(entry, schemas.SlideTextualContent):
                page = entry.slide
            else:
                page = 0

            # Page-level chunk (full text of the page/slide/section)
            chunks.append(
                models.ContentChunk(
                    content=entry.text,
                    page=page,
                    source_file=source_file,
                    chunk_index=chunk_index,
                    course_id=course_id,
                    chunk_type=models.ContentChunkType.PAGE,
                )
            )
            chunk_index += 1

            # Paragraph-level chunks (only if >1 non-empty paragraph)
            paragraphs = [p.strip() for p in entry.text.split("\n\n") if p.strip()]
            if len(paragraphs) > 1:
                for paragraph in paragraphs:
                    chunks.append(
                        models.ContentChunk(
                            content=paragraph,
                            page=page,
                            source_file=source_file,
                            chunk_index=chunk_index,
                            course_id=course_id,
                            chunk_type=models.ContentChunkType.PARAGRAPH,
                        )
                    )
                    chunk_index += 1

        self.__logger.info(f"Chunking complete. Total chunks: {len(chunks)}")
        return chunks
